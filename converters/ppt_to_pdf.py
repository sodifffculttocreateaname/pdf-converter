# -*- coding: utf-8 -*-
"""
PPT转PDF转换器
支持多种转换方式:
1. Windows COM (仅Windows，需要安装Microsoft PowerPoint)
2. python-pptx + reportlab (跨平台，但格式可能不完全准确)
"""
import platform
from pathlib import Path
from typing import List

from config.constants import ConversionStatus
from core.base_converter import BaseConverter, ConversionProgress, ConversionResult
from utils.file_utils import ensure_dir
from utils.logger import get_logger

logger = get_logger(__name__)


class PptToPdfConverter(BaseConverter):
    """PPT转PDF转换器"""

    name = "PPT转PDF"
    description = "将PowerPoint演示文稿转换为PDF"
    supported_input_formats = ["ppt", "pptx"]
    supported_output_formats = ["pdf"]

    def __init__(self):
        super().__init__()
        self._use_com = self._check_com_available()

    def _check_com_available(self) -> bool:
        """检查COM是否可用（仅Windows）"""
        if platform.system() != 'Windows':
            return False

        try:
            import win32com.client
            return True
        except ImportError:
            return False

    def convert(self, input_files: List[Path], output_dir: Path, **kwargs) -> List[ConversionResult]:
        """
        执行PPT转PDF转换

        Args:
            input_files: 输入PPT文件列表
            output_dir: 输出目录
            **kwargs: 额外参数

        Returns:
            转换结果列表
        """
        results = []

        # 确保输出目录存在
        ensure_dir(output_dir)

        # 选择转换方法
        if self._use_com:
            results = self._convert_with_com(input_files, output_dir)
        else:
            results = self._convert_with_pptx(input_files, output_dir)

        return results

    def _convert_with_com(self, input_files: List[Path], output_dir: Path) -> List[ConversionResult]:
        """
        使用Windows COM转换（需要安装Microsoft PowerPoint）

        Args:
            input_files: 输入文件列表
            output_dir: 输出目录

        Returns:
            转换结果列表
        """
        import win32com.client
        import pythoncom

        results = []
        ppt_app = None

        try:
            # 初始化COM
            pythoncom.CoInitialize()

            # 创建PowerPoint应用
            ppt_app = win32com.client.Dispatch("PowerPoint.Application")

            total_files = len(input_files)

            for i, input_file in enumerate(input_files):
                if self._check_cancelled():
                    break

                # 报告进度
                progress = ConversionProgress(
                    current=i + 1,
                    total=total_files,
                    current_file=input_file.name,
                    message=f"正在转换: {input_file.name}"
                )
                self._report_progress(progress)

                result = ConversionResult(input_file=input_file)

                try:
                    # 打开演示文稿
                    presentation = ppt_app.Presentations.Open(str(input_file.absolute()))

                    # 生成输出文件名
                    output_filename = self.get_output_filename(input_file, 'pdf')
                    output_path = output_dir / output_filename

                    # 导出为PDF
                    # ppSaveAsPDF = 32
                    presentation.SaveAs(str(output_path), FileFormat=32)
                    presentation.Close()

                    result.output_file = output_path
                    result.status = ConversionStatus.COMPLETED
                    result.message = "转换成功"

                    logger.info(f"PPT转PDF成功: {input_file.name}")

                except Exception as e:
                    result.status = ConversionStatus.FAILED
                    result.error = str(e)
                    logger.error(f"PPT转PDF失败: {input_file.name} - {e}")

                results.append(result)

        except Exception as e:
            logger.error(f"COM转换失败: {e}")
            return self._convert_with_pptx(input_files, output_dir)

        finally:
            # 关闭PowerPoint应用
            if ppt_app:
                try:
                    ppt_app.Quit()
                except Exception:
                    pass

            # 清理COM
            try:
                pythoncom.CoUninitialize()
            except Exception:
                pass

        # 报告完成进度
        if not self._check_cancelled():
            progress = ConversionProgress(
                current=len(input_files),
                total=len(input_files),
                message="处理完成"
            )
            self._report_progress(progress)

        return results

    def _convert_with_pptx(self, input_files: List[Path], output_dir: Path) -> List[ConversionResult]:
        """
        使用python-pptx和reportlab转换（跨平台，格式可能不完全准确）

        Args:
            input_files: 输入文件列表
            output_dir: 输出目录

        Returns:
            转换结果列表
        """
        from pptx import Presentation
        from reportlab.lib.pagesizes import landscape, A4
        from reportlab.lib.units import cm
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

        results = []

        # 注册字体
        self._register_font()

        total_files = len(input_files)

        for i, input_file in enumerate(input_files):
            if self._check_cancelled():
                break

            # 报告进度
            progress = ConversionProgress(
                current=i + 1,
                total=total_files,
                current_file=input_file.name,
                message=f"正在转换: {input_file.name}"
            )
            self._report_progress(progress)

            result = ConversionResult(input_file=input_file)

            try:
                # 读取PPT
                prs = Presentation(str(input_file))

                # 生成输出文件名
                output_filename = self.get_output_filename(input_file, 'pdf')
                output_path = output_dir / output_filename

                # 创建PDF (横向A4)
                pdf_doc = SimpleDocTemplate(
                    str(output_path),
                    pagesize=landscape(A4),
                    leftMargin=1.5*cm,
                    rightMargin=1.5*cm,
                    topMargin=1.5*cm,
                    bottomMargin=1.5*cm
                )

                # 创建样式
                styles = getSampleStyleSheet()
                title_style = ParagraphStyle(
                    'TitleStyle',
                    parent=styles['Title'],
                    fontName=self._font_name,
                    fontSize=24,
                    leading=30,
                )
                normal_style = ParagraphStyle(
                    'NormalStyle',
                    parent=styles['Normal'],
                    fontName=self._font_name,
                    fontSize=14,
                    leading=18,
                )

                # 构建内容
                story = []

                for slide_num, slide in enumerate(prs.slides, 1):
                    if self._check_cancelled():
                        break

                    # 添加幻灯片标题
                    story.append(Paragraph(f"幻灯片 {slide_num}", title_style))
                    story.append(Spacer(1, 12))

                    # 提取文本
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text = shape.text.strip()
                            if text:
                                # 转义特殊字符
                                text = text.replace('&', '&amp;')
                                text = text.replace('<', '&lt;')
                                text = text.replace('>', '&gt;')

                                for line in text.split('\n'):
                                    if line.strip():
                                        story.append(Paragraph(line.strip(), normal_style))

                    # 添加分页（除了最后一页）
                    if slide_num < len(prs.slides):
                        story.append(PageBreak())

                if self._check_cancelled():
                    result.status = ConversionStatus.CANCELLED
                    if output_path.exists():
                        output_path.unlink()
                    results.append(result)
                    continue

                # 生成PDF
                pdf_doc.build(story)

                result.output_file = output_path
                result.status = ConversionStatus.COMPLETED
                result.message = f"转换成功 ({len(prs.slides)} 张幻灯片)"

                logger.info(f"PPT转PDF成功: {input_file.name}")

            except Exception as e:
                result.status = ConversionStatus.FAILED
                result.error = str(e)
                logger.error(f"PPT转PDF失败: {input_file.name} - {e}")

            results.append(result)

        # 报告完成进度
        if not self._check_cancelled():
            progress = ConversionProgress(
                current=total_files,
                total=total_files,
                message="处理完成"
            )
            self._report_progress(progress)

        return results

    _font_registered = False
    _font_name = "SimSun"

    def _register_font(self):
        """注册中文字体"""
        if self._font_registered:
            return

        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont

        try:
            font_paths = [
                r"C:\Windows\Fonts\simhei.ttf",
                r"C:\Windows\Fonts\simsun.ttc",
                r"C:\Windows\Fonts\msyh.ttc",
                "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
                "/System/Library/Fonts/PingFang.ttc",
            ]

            for font_path in font_paths:
                if Path(font_path).exists():
                    try:
                        pdfmetrics.registerFont(TTFont(self._font_name, font_path))
                        self._font_registered = True
                        logger.info(f"成功注册字体: {font_path}")
                        return
                    except Exception:
                        continue

            self._font_name = "Helvetica"

        except Exception as e:
            logger.warning(f"注册字体失败: {e}")
            self._font_name = "Helvetica"

        self._font_registered = True

    def validate_input(self, file_path: Path) -> bool:
        """
        验证输入文件

        Args:
            file_path: 文件路径

        Returns:
            如果文件有效返回 True
        """
        if not file_path.exists():
            return False

        return file_path.suffix.lower() in ['.ppt', '.pptx']